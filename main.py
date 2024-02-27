import csv
import shutil
from pptx import Presentation

BASE_FILE='model.pptx'


def read_csv(file_name):
    with open(file_name, 'r', newline='', encoding='utf-8') as data:
        persons = csv.reader(data)
        data_persons = list(persons)
        attributes = data_persons[0]
        values = data_persons[1:]
        return [attributes,values]


def replace_value(nome_arquivo, texto_antigo, texto_novo):
    prs = Presentation(nome_arquivo)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.replace(texto_antigo, texto_novo)
    prs.save(nome_arquivo)


attributes,values = read_csv(file_name='data.csv')



for person in values:
    output_file = person[1]+'.pptx'
    shutil.copy(BASE_FILE, output_file)
    for index,item in enumerate(person):
        replace_value(output_file,"#-{attributes[index]}-#",item)
        print(attributes[index]+"="+str(item))
  
        