import csv
import shutil
from pptx import Presentation
import os

BASE_FILE='model.pptx'
OUT_PUT_DIR='outputs'

if not os.path.exists(OUT_PUT_DIR):
    os.makedirs(OUT_PUT_DIR)


def read_csv(file_name):
    with open(file_name, 'r', newline='', encoding='utf-8') as data:
        persons = csv.reader(data)
        data_persons = list(persons)
        attributes = data_persons[0]
        values = data_persons[1:]
        return [attributes,values]

def replace_value(nome_arquivo, texto_antigo, texto_novo):
    prs = Presentation(nome_arquivo)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape.text = shape.text.replace(texto_antigo, texto_novo)
    prs.save(nome_arquivo)


def do_replacement(input,output,dataset_file):
    attributes,values = read_csv(file_name=dataset_file)

    for person in values:
        output_file = output+"/"+person[1]+'.pptx'
        shutil.copy(input,output_file)
        for index,item in enumerate(person):
            replace_value(output_file,f"#-{attributes[index]}-#",item)
            print(attributes[index]+"="+str(item))
